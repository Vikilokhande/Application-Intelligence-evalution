import os
import sys
import time
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (Government Deep Navy & Warm Gold) ──────────────────────
COLOR_NAVY = RGBColor(10, 37, 64)          # #0A2540
COLOR_DARK_NAVY = RGBColor(7, 26, 43)      # #071A2B
COLOR_GOLD = RGBColor(197, 155, 39)        # #C59B27
COLOR_GOLD_LIGHT = RGBColor(255, 251, 235) # #FFFBEB
COLOR_SLATE_DARK = RGBColor(15, 23, 42)    # #0F172A
COLOR_SLATE_MUTED = RGBColor(71, 85, 105)  # #475569
COLOR_SLATE_LIGHT = RGBColor(241, 245, 249)# #F1F5F9
COLOR_BG_LIGHT = RGBColor(248, 250, 252)   # #F8FAFC
COLOR_WHITE = RGBColor(255, 255, 255)      # #FFFFFF
COLOR_GREEN = RGBColor(21, 128, 61)        # #15803D
COLOR_GREEN_BG = RGBColor(220, 252, 231)   # #DCFCE7
COLOR_RED = RGBColor(220, 38, 38)          # #DC2626
COLOR_RED_BG = RGBColor(254, 242, 242)     # #FEF2F2
COLOR_BORDER = RGBColor(226, 232, 240)     # #E2E8F0
COLOR_ARROW = RGBColor(148, 163, 184)      # #94A3B8

def build_presentation():
    project_root = Path(__file__).resolve().parents[1]
    output_pptx = project_root / "docs" / "DECC_Environmental_Review_Platform_Executive_Presentation.pptx"
    primary_pptx = project_root / "docs" / "DECC_Environmental_Review_Platform_Presentation.pptx"
    output_pptx.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="ENVIRONMENTAL APPLICATION REVIEW PORTAL"):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.95))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(10.5)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_GOLD
        p0.space_after = Pt(2)

        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_NAVY

    def add_card(slide, left, top, width, height, bg_color=COLOR_WHITE, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2)
        return shape

    def add_arrow_badge(slide, left, top, width=Inches(0.4), height=Inches(0.4)):
        shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_GOLD
        shape.line.fill.background()
        return shape

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE SLIDE (Clean & Human-Friendly)
    # ══════════════════════════════════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_NAVY
    bg1.line.fill.background()

    gold_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.2), Inches(1.2), Inches(0.08))
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = COLOR_GOLD
    gold_bar.line.fill.background()

    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.933), Inches(4.8))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "GOVERNMENT OF MAHARASHTRA • DIRECTORATE OF ENVIRONMENT & CLIMATE CHANGE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(14)

    p = tf1.add_paragraph()
    p.text = "Smart Environmental Application Review Platform"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(14)

    p = tf1.add_paragraph()
    p.text = "An easy-to-use digital system that reads citizen applications, cross-checks documents for errors, verifies government scheme rules, and calculates a clear risk score to help officers make fast, fair decisions."
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.space_after = Pt(28)

    p = tf1.add_paragraph()
    p.text = "CORE PRINCIPLE:  AI ASSISTS · THE OFFICER DECIDES"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: ACTUAL PROBLEM STATEMENT (Human-Friendly)
    # ══════════════════════════════════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "The Real-World Problem: Why Manual Review Is Slow & Difficult", "THE CHALLENGE")

    cards_data_2 = [
        ("⏳ 15 to 30 Days of Manual Waiting", "Officers must read dozens of physical pages and scanned PDFs line-by-line. This causes long citizen waiting times and huge office backlogs.", COLOR_RED_BG, COLOR_RED),
        ("🔍 Hard to Spot Hidden Contradictions", "If the project cost is ₹45 Lakhs in the proposal but ₹50 Lakhs in the budget sheet, catching these small differences across multiple files is very difficult by eye.", COLOR_RED_BG, COLOR_RED),
        ("📚 Complex Rules Scattered in PDFs", "Government scheme guidelines, budget caps, and eligibility limits are buried in long policy PDFs, making it hard to double-check every single rule.", COLOR_RED_BG, COLOR_RED),
        ("⚖️ No Clear Way to Prioritize Cases", "Officers have no quick risk summary. Clean, straightforward applications take just as long to process as suspicious or high-risk ones.", COLOR_RED_BG, COLOR_RED),
    ]

    for i, (title, desc, bg, accent) in enumerate(cards_data_2):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.5 + row * 2.7)
        add_card(slide2, left, top, Inches(5.75), Inches(2.45), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
        
        strip = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.15), Inches(2.45))
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent
        strip.line.fill.background()

        tb = slide2.shapes.add_textbox(left + Inches(0.35), top + Inches(0.25), Inches(5.15), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_SLATE_MUTED

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: OUR SOLUTION (Clear & Practical)
    # ══════════════════════════════════════════════════════════════════════════
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Our Solution: Fast Automated Checks with Complete Officer Control", "WHAT WE BUILT")

    cards_data_3 = [
        ("Smart Document Reader", "Reads digital PDFs, scanned files, spreadsheets, and certificates automatically using smart OCR technology.", "STEP 1: READ"),
        ("Cross-Document Checker", "Compares all files at once to instantly spot any mismatches in project costs, dates, or applicant names.", "STEP 2: COMPARE"),
        ("Instant Policy Matcher", "Searches government policy guides to make sure the project fits all legal guidelines and funding limits.", "STEP 3: MATCH"),
        ("Easy 0–100 Risk Score", "Uses machine learning to give each application a clear Risk Score (0-100) and Quality Score (0-100).", "STEP 4: SCORE"),
        ("Clear Plain-English Summary", "Provides a 2-minute case summary and suggested questions for the applicant without any confusing jargon.", "STEP 5: EXPLAIN"),
        ("Officer Decision Cockpit", "The officer makes the final call (Approve, Clarify, or Reject) and sends a clear email report to the citizen in 1 click.", "STEP 6: DECIDE"),
    ]

    for i, (title, desc, badge) in enumerate(cards_data_3):
        row = i // 3
        col = i % 3
        left = Inches(0.8 + col * 3.95)
        top = Inches(1.5 + row * 2.7)
        add_card(slide3, left, top, Inches(3.75), Inches(2.45), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)

        tb = slide3.shapes.add_textbox(left + Inches(0.25), top + Inches(0.25), Inches(3.25), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = badge
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SLATE_MUTED

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4: HIGH-LEVEL ARCHITECTURE DIAGRAM (Simple & Visual)
    # ══════════════════════════════════════════════════════════════════════════
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "How the System Works: Simple High-Level Architecture", "SYSTEM ARCHITECTURE")

    # Column 1: Citizen Input (Left)
    col1_left = Inches(0.8)
    col1_width = Inches(2.5)
    add_card(slide4, col1_left, Inches(1.5), col1_width, Inches(5.2), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
    
    tb1 = slide4.shapes.add_textbox(col1_left + Inches(0.15), Inches(1.7), col1_width - Inches(0.3), Inches(4.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
    p = tf1.paragraphs[0]
    p.text = "1. CITIZEN INPUT"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(8)
    p = tf1.add_paragraph()
    p.text = "Citizen / Organization"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    p = tf1.add_paragraph()
    p.text = "Submits online application with attached files:\n• Project Proposal (PDF)\n• Budget Sheet (XLSX)\n• EIA Environmental Report\n• Organization Certificate"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_SLATE_MUTED

    # Arrow 1 -> Center
    add_arrow_badge(slide4, Inches(3.45), Inches(3.9), Inches(0.4), Inches(0.35))

    # Column 2: Smart Processing Engines (Center Box)
    col2_left = Inches(4.0)
    col2_width = Inches(5.333)
    add_card(slide4, col2_left, Inches(1.5), col2_width, Inches(5.2), bg_color=COLOR_SLATE_LIGHT, border_color=COLOR_BORDER)
    
    tb2_hdr = slide4.shapes.add_textbox(col2_left + Inches(0.2), Inches(1.7), col2_width - Inches(0.4), Inches(0.5))
    tf2_h = tb2_hdr.text_frame
    tf2_h.margin_left = tf2_h.margin_top = tf2_h.margin_right = tf2_h.margin_bottom = 0
    p = tf2_h.paragraphs[0]
    p.text = "2. AUTOMATED VERIFICATION ENGINES"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD

    # 4 Inner Blocks in Center
    inner_blocks = [
        ("Smart OCR & Document Reader", "Converts scanned pages into clean digital data.", Inches(2.2)),
        ("Cross-Document Mismatch Checker", "Spots differences between budget and proposal.", Inches(3.2)),
        ("Scheme Policy & Guidelines Matcher", "Checks limits in official government guides.", Inches(4.2)),
        ("AI Risk & Quality Scoring Engine", "Calculates clear 0-100 risk & quality scores.", Inches(5.2)),
    ]
    for b_title, b_desc, b_top in inner_blocks:
        add_card(slide4, col2_left + Inches(0.2), b_top, col2_width - Inches(0.4), Inches(0.85), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
        tbi = slide4.shapes.add_textbox(col2_left + Inches(0.35), b_top + Inches(0.1), col2_width - Inches(0.7), Inches(0.65))
        tfi = tbi.text_frame
        tfi.margin_left = tfi.margin_top = tfi.margin_right = tfi.margin_bottom = 0
        p = tfi.paragraphs[0]
        p.text = b_title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p = tfi.add_paragraph()
        p.text = b_desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_SLATE_MUTED

    # Arrow 2 -> Right
    add_arrow_badge(slide4, Inches(9.48), Inches(3.9), Inches(0.4), Inches(0.35))

    # Column 3: Officer Review & Outcome (Right)
    col3_left = Inches(10.033)
    col3_width = Inches(2.5)
    add_card(slide4, col3_left, Inches(1.5), col3_width, Inches(5.2), bg_color=COLOR_GOLD_LIGHT, border_color=COLOR_GOLD)
    
    tb3 = slide4.shapes.add_textbox(col3_left + Inches(0.15), Inches(1.7), col3_width - Inches(0.3), Inches(4.8))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
    p = tf3.paragraphs[0]
    p.text = "3. OFFICER DECISION"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(8)
    p = tf3.add_paragraph()
    p.text = "Government Officer"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    p = tf3.add_paragraph()
    p.text = "1. Reviews verified evidence\n2. Checks risk score\n3. Makes final decision:\n   • Approve\n   • Request Clarification\n   • Reject\n4. Dispatches official email clearance report in 1 click"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_SLATE_MUTED

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5: END-TO-END WORKFLOW (Simple 7 Steps with Arrows)
    # ══════════════════════════════════════════════════════════════════════════
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "End-to-End Journey: 7 Simple Steps from Upload to Clearance", "STEP-BY-STEP WORKFLOW")

    steps_7 = [
        ("01", "Application Submitted", "Citizen fills details & uploads files."),
        ("02", "Files Checked", "System verifies file formats & safety."),
        ("03", "Data Read (OCR)", "Smart reader extracts text & tables."),
        ("04", "Cross-Checked", "Compares numbers across all files."),
        ("05", "Rules Verified", "Checks government scheme limits."),
        ("06", "AI Risk Scored", "Computes clear 0-100 risk score."),
        ("07", "Officer Decision", "Officer signs off & emails citizen report."),
    ]

    for i, (num, title, desc) in enumerate(steps_7):
        left = Inches(0.8 + i * 1.69)
        top = Inches(1.8)
        width = Inches(1.55)
        height = Inches(4.6)

        is_last = (i == 6)
        card_bg = COLOR_GOLD_LIGHT if is_last else COLOR_WHITE
        border = COLOR_GOLD if is_last else COLOR_BORDER

        add_card(slide5, left, top, width, height, bg_color=card_bg, border_color=border)

        tb = slide5.shapes.add_textbox(left + Inches(0.1), top + Inches(0.2), width - Inches(0.2), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = f"STEP {num}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD if is_last else COLOR_NAVY
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p.space_after = Pt(10)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10.5)
        p.font.color.rgb = COLOR_SLATE_MUTED

        # Small arrow between cards (except last)
        if i < 6:
            arrow_left = left + width + Inches(0.02)
            add_arrow_badge(slide5, arrow_left, top + Inches(1.8), Inches(0.12), Inches(0.2))

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6: 5 UNIQUE SUPERPOWERS (Human-Friendly Features)
    # ══════════════════════════════════════════════════════════════════════════
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Why Our Platform Is Special: 5 Key Superpowers", "KEY INNOVATIONS")

    features_6 = [
        ("1. 100% Reliable Legal Checks (No Guesswork)", "Mandatory fields, budget limits, and scheme eligibility are calculated using exact mathematical rules. The AI never guesses or makes up fake pass results."),
        ("2. Automatic Conflict Spotter", "If the proposal mentions a 12-month timeline but the schedule says 24 months, the system immediately flags this contradiction for the officer to review."),
        ("3. Instant Policy Clause Finder", "The system searches official government scheme documents and highlights the exact policy paragraph relevant to the application."),
        ("4. Objective 0–100 Risk & Quality Index", "Instead of a vague recommendation, officers see clear numerical scores (e.g. Risk: 15/100, Quality: 92/100) showing exactly which factors contributed."),
        ("5. Always Stays Working (Resilient Design)", "Even if internet AI services experience temporary delays, all core checks, document readers, and scoring models continue working smoothly without stopping."),
    ]

    for i, (title, desc) in enumerate(features_6):
        top = Inches(1.5 + i * 1.1)
        add_card(slide6, Inches(0.8), top, Inches(11.733), Inches(0.95), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)

        tb = slide6.shapes.add_textbox(Inches(1.1), top + Inches(0.12), Inches(11.133), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SLATE_MUTED

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7: REVIEWER COCKPIT & EMAIL (Human-Centered)
    # ══════════════════════════════════════════════════════════════════════════
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Officer Cockpit & Instant Citizen Communication", "DECISION & NOTIFICATION")

    # Left: Officer Workspace
    add_card(slide7, Inches(0.8), Inches(1.5), Inches(5.75), Inches(5.2), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
    tb = slide7.shapes.add_textbox(Inches(1.1), Inches(1.75), Inches(5.15), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "FOR GOVERNMENT OFFICERS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = "The Officer Decision Cockpit"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(10)

    bullets_l = [
        ("Clear Step-by-Step Overview", "View application details, flagged mismatches, and scheme rules all in one clean screen."),
        ("One-Click Decisions", "Choose Approve, Request Clarification, or Reject based on verified facts."),
        ("Full Legal Accountability", "If an officer chooses to override the AI advice, their written reason is securely recorded for government audit."),
        ("English & Hindi Language Support", "Switch between English and Hindi (मराठी/हिंदी) instantly for total ease of use."),
    ]
    for b_title, b_desc in bullets_l:
        p = tf.add_paragraph()
        p.text = f"•  {b_title}: {b_desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SLATE_MUTED
        p.space_after = Pt(8)

    # Right: Email Reports
    add_card(slide7, Inches(6.78), Inches(1.5), Inches(5.75), Inches(5.2), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
    tb = slide7.shapes.add_textbox(Inches(7.08), Inches(1.75), Inches(5.15), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "FOR CITIZENS & APPLICANTS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = "Clear Email Clearance Reports"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(10)

    bullets_r = [
        ("1-Click Email Dispatch", "Officers click 'Email Report' to send an official digital summary directly to the citizen's inbox."),
        ("Simple Action Cards", "If documents are missing or numbers don't match, the email clearly explains what the citizen needs to fix."),
        ("Clean Mobile & Desktop Layout", "Emails look beautiful on both mobile phones and desktop computers with official government branding."),
        ("Instant Delivery", "Sent securely in seconds via standard email relays."),
    ]
    for b_title, b_desc in bullets_r:
        p = tf.add_paragraph()
        p.text = f"•  {b_title}: {b_desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SLATE_MUTED
        p.space_after = Pt(8)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8: IMPACT & SUMMARY (Inspiring & Clear)
    # ══════════════════════════════════════════════════════════════════════════
    slide8 = prs.slides.add_slide(blank_layout)
    bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = COLOR_NAVY
    bg8.line.fill.background()

    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(5.8))
    tf8 = tb.text_frame
    tf8.word_wrap = True

    p = tf8.paragraphs[0]
    p.text = "REAL-WORLD IMPACT"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(10)

    p = tf8.add_paragraph()
    p.text = "Faster Reviews. Total Transparency. Zero Guesswork."
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(20)

    metrics = [
        ("90% Faster for Citizens", "Citizens get responses in days instead of waiting up to a month for manual file review."),
        ("100% Transparent & Auditable", "Every single calculation, rule test, and officer decision is permanently saved for total accountability."),
        ("Zero Mistakes in Policy Rules", "Mathematical checks and policy guideline search eliminate accidental human oversights."),
    ]

    for m_title, m_desc in metrics:
        p = tf8.add_paragraph()
        p.text = f"✓  {m_title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD
        p.space_after = Pt(2)

        p = tf8.add_paragraph()
        p.text = f"    {m_desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.space_after = Pt(14)

    # Save handling
    saved = False
    for target_path in [primary_pptx, output_pptx]:
        try:
            prs.save(str(target_path))
            print(f"Successfully saved PowerPoint presentation: {target_path}")
            saved = True
        except Exception as e:
            print(f"Note on saving {target_path.name}: {e}")

    if not saved:
        timestamp_path = project_root / "docs" / f"DECC_Review_Platform_Presentation_{int(time.time())}.pptx"
        prs.save(str(timestamp_path))
        print(f"Successfully saved presentation to alternative path: {timestamp_path}")

if __name__ == "__main__":
    build_presentation()
